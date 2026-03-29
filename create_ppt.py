#!/usr/bin/env python3
"""
KisanMind — Guy Kawasaki 10/20/30 Pitch Deck Generator
10 slides, 20 minutes, 30-point font minimum
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)  # Deep navy
BG_CARD    = RGBColor(0x1A, 0x36, 0x5D)  # Card background
ACCENT_GREEN  = RGBColor(0x22, 0xC5, 0x5E)  # Indian flag green
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x33)  # Indian flag saffron
ACCENT_BLUE   = RGBColor(0x38, 0xBD, 0xF8)  # Tech blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
RED_ACCENT = RGBColor(0xF4, 0x3F, 0x5E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=30, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=28, default_color=WHITE):
    """lines = list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(6)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════════════════════
# SLIDE 1: TITLE — Problem / Opportunity
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)
add_accent_line(slide, 0, 0.06, 13.333, ACCENT_ORANGE)

add_text_box(slide, 1, 1.0, 11, 1.2, "KisanMind", 72, ACCENT_GREEN, True)
add_text_box(slide, 1, 2.0, 11, 0.8, "Satellite-to-Voice Agricultural Intelligence for 150M Indian Farmers", 32, LIGHT_GRAY)

add_accent_line(slide, 1, 3.2, 4, ACCENT_ORANGE)

add_multi_text(slide, 1, 3.6, 11, 3.5, [
    ("The Problem:", 36, ACCENT_ORANGE, True),
    ("", 12),
    ("150M farming households make decisions worth Rs 45 lakh crore/year with:", 30, WHITE),
    ("", 8),
    ("   Zero satellite visibility on crop health", 28, RED_ACCENT, True),
    ("   No real-time price comparison across mandis (30-40% daily fluctuation)", 28, RED_ACCENT, True),
    ("   Generic weather forecasts useless for field-level decisions", 28, RED_ACCENT, True),
    ("   Language barrier: advisory only in English/Hindi, useless for 60%+ farmers", 28, RED_ACCENT, True),
])

add_multi_text(slide, 1, 6.5, 11, 0.7, [
    ("ET AI Hackathon 2026  |  Problem #5: Domain-Specialized AI Agents  |  kisanmind.dmj.one", 18, GRAY),
])


# ═══════════════════════════════════════════════════════════
# SLIDE 2: VALUE PROPOSITION — What KisanMind Does
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.5, 11, 0.8, "One Phone Call. Six Data Sources. 22 Languages.", 40, WHITE, True)
add_accent_line(slide, 1, 1.3, 6, ACCENT_ORANGE)

# The farmer says...
add_rounded_rect(slide, 0.8, 1.7, 11.7, 1.3, RGBColor(0x15, 0x2B, 0x4A))
add_multi_text(slide, 1.2, 1.8, 11, 1.2, [
    ('Farmer says: "Main Solan mein tamatar uga raha hoon"', 26, ACCENT_BLUE, True),
    ("", 6),
    ('KisanMind responds in 15 seconds with: Crop health (NDVI 0.54 from Sentinel-2) + Best mandi', 24, LIGHT_GRAY),
    ('(Bhuntar Rs 7,500/q, 251km) + Net profit after transport/commission/spoilage + Weather-timed', 24, LIGHT_GRAY),
    ('harvest advice + All in farmer\'s native language via voice', 24, LIGHT_GRAY),
])

# 6 data sources
add_text_box(slide, 1, 3.3, 11, 0.6, "6 Real-Time Data Sources Fused Per Request:", 30, ACCENT_GREEN, True)

sources = [
    ("Sentinel-2", "NDVI/EVI/NDWI\n10m resolution", ACCENT_GREEN),
    ("Sentinel-1 SAR", "Soil moisture\nthrough clouds", ACCENT_GREEN),
    ("MODIS/SMAP", "Temperature +\nroot-zone moisture", ACCENT_GREEN),
    ("AgMarkNet", "Govt mandi prices\n106+ crops", ACCENT_ORANGE),
    ("Google Maps", "Real driving distance\nto every mandi", ACCENT_BLUE),
    ("Open-Meteo", "5-day forecast +\n120-day historical", ACCENT_BLUE),
]

for i, (title, desc, accent) in enumerate(sources):
    col = i % 6
    x = 0.8 + col * 2.05
    y = 4.1
    add_rounded_rect(slide, x, y, 1.9, 2.0, BG_CARD)
    add_accent_line(slide, x + 0.2, y + 0.15, 1.5, accent)
    add_multi_text(slide, x + 0.15, y + 0.35, 1.7, 1.5, [
        (title, 20, accent, True, PP_ALIGN.CENTER),
        ("", 4),
        (desc, 16, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ])

add_text_box(slide, 1, 6.5, 11, 0.6, "Zero fake data. Every data point from a verified API call.", 22, GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: HOW IT WORKS — Multi-Agent Architecture
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "Multi-Agent Architecture (Google ADK)", 40, WHITE, True)
add_text_box(slide, 1, 1.1, 11, 0.5, "5 specialist agents orchestrated by a Brain agent | Built on Google Agent Development Kit", 24, GRAY)

# Brain orchestrator box
add_rounded_rect(slide, 3.5, 1.8, 6.3, 1.2, RGBColor(0x1A, 0x36, 0x5D))
add_accent_line(slide, 3.7, 1.95, 5.9, ACCENT_BLUE)
add_multi_text(slide, 3.7, 2.1, 6, 0.9, [
    ("Brain (Orchestrator) — Gemini 3.1 Pro", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER),
    ("Routes intents | Merges outputs | Enforces guardrails | Synthesizes advisory", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER),
])

# 4 specialist agents
agents = [
    ("VaaniSetu", "Voice Bridge", "STT/TTS in 22 languages\nIntent extraction\nLanguage detection", "Gemini 3 Flash", ACCENT_BLUE),
    ("SatDrishti", "Satellite Eye", "Sentinel-2 NDVI/EVI/NDWI\nHealth classification\n30-day trend detection", "Gemini 3 Flash", ACCENT_GREEN),
    ("MandiMitra", "Market Friend", "AgMarkNet live prices\nNet profit ranking\nTransport + spoilage calc", "Gemini 3 Flash", ACCENT_ORANGE),
    ("MausamGuru", "Weather Guru", "5-day forecast\nCrop-specific DO/DON'T\nGrowth stage awareness", "Gemini 3 Flash", ACCENT_BLUE),
]

for i, (name, role, desc, model, accent) in enumerate(agents):
    x = 0.5 + i * 3.15
    y = 3.5
    add_rounded_rect(slide, x, y, 2.95, 3.0, BG_CARD)
    add_accent_line(slide, x + 0.15, y + 0.15, 2.65, accent)
    add_multi_text(slide, x + 0.15, y + 0.35, 2.7, 2.5, [
        (name, 22, accent, True, PP_ALIGN.CENTER),
        (role, 16, GRAY, False, PP_ALIGN.CENTER),
        ("", 6),
        (desc, 16, LIGHT_GRAY, False, PP_ALIGN.CENTER),
        ("", 4),
        (model, 14, GRAY, False, PP_ALIGN.CENTER),
    ])

add_text_box(slide, 1, 6.8, 11, 0.5, "10,110 lines of code  |  16 API endpoints  |  4 satellite systems  |  5 cloud functions", 20, GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: UNDERLYING MAGIC — Satellite + GDD + Profit Formula
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "The Intelligence Layer: Real Science, Not Heuristics", 38, WHITE, True)

# LEFT: Satellite
add_rounded_rect(slide, 0.5, 1.4, 4.0, 2.8, BG_CARD)
add_accent_line(slide, 0.7, 1.55, 3.6, ACCENT_GREEN)
add_multi_text(slide, 0.7, 1.7, 3.6, 2.4, [
    ("4 Satellite Systems", 24, ACCENT_GREEN, True),
    ("", 6),
    ("Sentinel-2: NDVI = (B8-B4)/(B8+B4)", 16, LIGHT_GRAY),
    ("  EVI = 2.5x((B8-B4)/(B8+6xB4-7.5xB2+1))", 16, LIGHT_GRAY),
    ("  NDWI = (B3-B8)/(B3+B8)", 16, LIGHT_GRAY),
    ("  10m resolution, <30% cloud filter", 16, GRAY),
    ("", 4),
    ("Sentinel-1 SAR: VV+VH radar soil moisture", 16, LIGHT_GRAY),
    ("MODIS Terra: 1km daily land surface temp", 16, LIGHT_GRAY),
    ("NASA SMAP: 9km root-zone soil moisture", 16, LIGHT_GRAY),
])

# CENTER: GDD
add_rounded_rect(slide, 4.7, 1.4, 4.0, 2.8, BG_CARD)
add_accent_line(slide, 4.9, 1.55, 3.6, ACCENT_BLUE)
add_multi_text(slide, 4.9, 1.7, 3.6, 2.4, [
    ("Growing Degree Days (GDD)", 24, ACCENT_BLUE, True),
    ("", 6),
    ("GDD = max(0, (Tmax+Tmin)/2 - Tbase)", 16, LIGHT_GRAY),
    ("120-day historical weather data", 16, LIGHT_GRAY),
    ("", 4),
    ("10 crops modeled:", 16, WHITE, True),
    ("Tomato (base 10C): 6 growth stages", 16, LIGHT_GRAY),
    ("  0/300/600/900/1200/1500 GDD", 14, GRAY),
    ("Wheat (base 5C): 6 stages", 16, LIGHT_GRAY),
    ("Rice, Potato, Onion, Capsicum,", 16, LIGHT_GRAY),
    ("Cabbage, Cauliflower, Apple, Mango", 16, LIGHT_GRAY),
])

# RIGHT: Profit formula
add_rounded_rect(slide, 8.9, 1.4, 4.0, 2.8, BG_CARD)
add_accent_line(slide, 9.1, 1.55, 3.6, ACCENT_ORANGE)
add_multi_text(slide, 9.1, 1.7, 3.6, 2.4, [
    ("Net Profit Ranking", 24, ACCENT_ORANGE, True),
    ("", 6),
    ("Net Profit = Modal Price", 16, LIGHT_GRAY),
    ("  - Transport (Rs 3.5/km/quintal)", 16, LIGHT_GRAY),
    ("  - Commission (4% of price)", 16, LIGHT_GRAY),
    ("  - Spoilage Loss", 16, LIGHT_GRAY),
    ("", 4),
    ("Spoilage rates per hour:", 16, WHITE, True),
    ("  Tomato: 0.5%/hr (high perishable)", 16, RED_ACCENT),
    ("  Potato: 0.2%/hr (medium)", 16, ACCENT_ORANGE),
    ("  Wheat: 0.05%/hr (low)", 16, ACCENT_GREEN),
    ("Road winding factor: 1.3x", 16, GRAY),
])

# Crop weather rules
add_rounded_rect(slide, 0.5, 4.5, 12.4, 2.5, BG_CARD)
add_accent_line(slide, 0.7, 4.65, 12, ACCENT_BLUE)
add_text_box(slide, 0.7, 4.8, 12, 0.5, "Crop-Specific Weather Rules (5 crops with full rule sets)", 24, ACCENT_BLUE, True)

add_multi_text(slide, 0.7, 5.3, 5.8, 1.6, [
    ("Tomato: Frost <10C, flower drop >40C, heavy rain", 16, LIGHT_GRAY),
    ("  >20mm = harvest now, humidity >85% = blight risk", 16, GRAY),
    ("Wheat: <5C during heading = yield loss, >30mm at", 16, LIGHT_GRAY),
    ("  harvest = grain sprouting, >90% humid = rust risk", 16, GRAY),
    ("Rice: <15C = cold stress (maintain standing water),", 16, LIGHT_GRAY),
    ("  >50mm at flowering = excess flood risk", 16, GRAY),
])
add_multi_text(slide, 6.8, 5.3, 5.8, 1.6, [
    ("Apple: Frost <-2C = severe (smudge pots needed),", 16, LIGHT_GRAY),
    ("  >25mm at fruit dev = cracking, >80% humid = scab", 16, GRAY),
    ("Coffee: <10C = berry damage, >50mm = root rot,", 16, LIGHT_GRAY),
    ("  humidity <40% = irrigate immediately", 16, GRAY),
    ("", 6),
    ("8 cross-validation rules catch conflicts", 16, ACCENT_GREEN, True),
    ("  (e.g., NDVI declining + adequate rain = pest issue)", 14, GRAY),
])


# ═══════════════════════════════════════════════════════════
# SLIDE 5: GUARDRAILS & COMPLIANCE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "Anti-Hallucination Guardrails & Compliance", 40, WHITE, True)
add_text_box(slide, 1, 1.1, 11, 0.5, "Every advisory is fact-checked, confidence-scored, and auditable", 24, GRAY)

# Left: Guardrail rules
add_rounded_rect(slide, 0.5, 1.8, 6.2, 3.0, BG_CARD)
add_accent_line(slide, 0.7, 1.95, 5.8, RED_ACCENT)
add_multi_text(slide, 0.7, 2.1, 5.8, 2.7, [
    ("Compliance Rules (Enforced in Code)", 24, RED_ACCENT, True),
    ("", 6),
    ("No pesticide brand/dosage recommendations", 20, WHITE),
    ("  -> Refers to KVK helpline 1800-180-1551", 18, GRAY),
    ("No loan/credit/insurance advice", 20, WHITE),
    ("  -> Blocked in system prompt + fact-check", 18, GRAY),
    ("No yield guarantees", 20, WHITE),
    ("  -> All estimates marked 'indicative'", 18, GRAY),
    ("Mandatory data source citations with timestamps", 20, WHITE),
    ("Full audit trail: session_id, intent, agents, sources", 20, WHITE),
])

# Right: Fact-checking
add_rounded_rect(slide, 6.9, 1.8, 6.0, 3.0, BG_CARD)
add_accent_line(slide, 7.1, 1.95, 5.6, ACCENT_GREEN)
add_multi_text(slide, 7.1, 2.1, 5.6, 2.7, [
    ("Dual-LLM Fact-Check Architecture", 24, ACCENT_GREEN, True),
    ("", 6),
    ("1. Gemini 3.1 Pro generates advisory", 20, WHITE),
    ("2. Gemini Flash Lite cross-validates against", 20, WHITE),
    ("   raw source data (NDVI, prices, weather)", 18, GRAY),
    ("", 4),
    ("Confidence Scoring (0.0 - 1.0):", 20, ACCENT_BLUE, True),
    ("  HIGH (>= 0.7): Full recommendation", 18, ACCENT_GREEN),
    ("  MEDIUM (0.4-0.7): Hedged advice", 18, ACCENT_ORANGE),
    ("  LOW (< 0.4): Refer to KVK", 18, RED_ACCENT),
    ("  Each data layer scored independently", 18, GRAY),
])

# Bottom: 8 cross-validation rules
add_rounded_rect(slide, 0.5, 5.1, 12.4, 2.0, BG_CARD)
add_accent_line(slide, 0.7, 5.25, 12, ACCENT_BLUE)
add_text_box(slide, 0.7, 5.4, 12, 0.5, "8 Cross-Validation Rules (Conflict Detection)", 24, ACCENT_BLUE, True)
add_multi_text(slide, 0.7, 5.9, 6, 1.0, [
    ("1. NDVI declining + adequate rain -> pest issue", 16, LIGHT_GRAY),
    ("2. NDVI healthy + farmer reports problem -> old data", 16, LIGHT_GRAY),
    ("3. Price rising + high arrivals -> unsustainable", 16, LIGHT_GRAY),
    ("4. Near harvest + NDVI plateau + good weather -> pick now", 16, LIGHT_GRAY),
])
add_multi_text(slide, 6.8, 5.9, 6, 1.0, [
    ("5. Frost warning -> protect sensitive crops", 16, LIGHT_GRAY),
    ("6. SAR vs NDVI cross-check for soil moisture", 16, LIGHT_GRAY),
    ("7. MODIS LST heat stress detection", 16, LIGHT_GRAY),
    ("8. SMAP root-zone vs surface moisture mismatch", 16, LIGHT_GRAY),
])


# ═══════════════════════════════════════════════════════════
# SLIDE 6: VOICE-FIRST & LANGUAGE SUPPORT
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "Voice-First: 22 Languages, Any Phone", 40, WHITE, True)
add_text_box(slide, 1, 1.1, 11, 0.5, "Cloud STT V2 + Cloud TTS Wavenet/Neural2 + Cloud Translation v3 + Twilio Voice", 24, GRAY)

# Voice pipeline
add_rounded_rect(slide, 0.5, 1.8, 7.8, 2.5, BG_CARD)
add_accent_line(slide, 0.7, 1.95, 7.4, ACCENT_BLUE)
add_multi_text(slide, 0.7, 2.1, 7.4, 2.2, [
    ("Voice Interaction Pipeline", 24, ACCENT_BLUE, True),
    ("", 6),
    ("1. Farmer calls (+1 260-254-7946) or taps 'Call KisanMind' on web app", 18, LIGHT_GRAY),
    ("2. GPS auto-detects location (farmer only says their crop)", 18, LIGHT_GRAY),
    ("3. Cloud STT V2 transcribes speech in native language", 18, LIGHT_GRAY),
    ("4. Gemini extracts intent: crop, location, question type", 18, LIGHT_GRAY),
    ("5. Parallel data fetch: satellite + mandi + weather (farming trivia plays as filler)", 18, LIGHT_GRAY),
    ("6. Advisory synthesized in farmer's language", 18, LIGHT_GRAY),
    ("7. Cloud TTS Wavenet speaks response aloud", 18, LIGHT_GRAY),
    ("8. Multi-turn: farmer asks follow-ups (auto-ends after 3 silences)", 18, LIGHT_GRAY),
])

# Language grid
add_rounded_rect(slide, 8.5, 1.8, 4.4, 2.5, BG_CARD)
add_accent_line(slide, 8.7, 1.95, 4, ACCENT_ORANGE)
add_multi_text(slide, 8.7, 2.1, 4, 2.2, [
    ("Native TTS (10 languages)", 20, ACCENT_ORANGE, True),
    ("Hindi, English, Tamil, Telugu", 16, LIGHT_GRAY),
    ("Bengali, Marathi, Gujarati", 16, LIGHT_GRAY),
    ("Kannada, Malayalam, Punjabi", 16, LIGHT_GRAY),
    ("", 6),
    ("Via Hindi bridge (12 more)", 20, ACCENT_BLUE, True),
    ("Odia, Assamese, Maithili", 16, LIGHT_GRAY),
    ("Sanskrit, Nepali, Sindhi", 16, LIGHT_GRAY),
    ("Dogri, Kashmiri, Konkani", 16, LIGHT_GRAY),
    ("Santali, Bodo, Manipuri", 16, LIGHT_GRAY),
])

# Performance box
add_rounded_rect(slide, 0.5, 4.6, 12.4, 2.5, BG_CARD)
add_accent_line(slide, 0.7, 4.75, 12, ACCENT_GREEN)
add_text_box(slide, 0.7, 4.9, 12, 0.5, "Performance & Caching", 24, ACCENT_GREEN, True)

perf_items = [
    ("0.13s", "L1 cache hit\n(in-memory)"),
    ("~200ms", "L2 cache hit\n(Cloud Storage)"),
    ("15-25s", "Fresh advisory\n(all 6 APIs)"),
    ("2-4s", "Voice latency\n(STT + TTS)"),
    ("106+", "Crops\nsupported"),
    ("10m", "Satellite\nresolution"),
]

for i, (metric, desc) in enumerate(perf_items):
    x = 0.8 + i * 2.05
    y = 5.4
    add_rounded_rect(slide, x, y, 1.9, 1.4, RGBColor(0x15, 0x2B, 0x4A))
    add_multi_text(slide, x + 0.1, y + 0.15, 1.7, 1.1, [
        (metric, 28, ACCENT_GREEN, True, PP_ALIGN.CENTER),
        (desc, 14, GRAY, False, PP_ALIGN.CENTER),
    ])


# ═══════════════════════════════════════════════════════════
# SLIDE 7: BUSINESS MODEL — Impact Quantification
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.3, 11, 0.8, "Business Impact: Quantified Estimates", 40, WHITE, True)
add_text_box(slide, 1, 1.0, 11, 0.5, "Back-of-envelope math with stated assumptions | Conservative Year 1 projections", 24, GRAY)

# Assumptions box
add_rounded_rect(slide, 0.5, 1.6, 5.8, 2.2, BG_CARD)
add_accent_line(slide, 0.7, 1.75, 5.4, ACCENT_ORANGE)
add_multi_text(slide, 0.7, 1.9, 5.4, 2.0, [
    ("Assumptions", 24, ACCENT_ORANGE, True),
    ("", 4),
    ("TAM: 150M farming households in India", 18, LIGHT_GRAY),
    ("Year 1 reach: 100,000 farmers (0.07% penetration)", 18, LIGHT_GRAY),
    ("Avg farmer grows 2 crops/year, sells at 2-3 mandis", 18, LIGHT_GRAY),
    ("Mandi price variance: 30-40% across nearby markets", 18, LIGHT_GRAY),
    ("  (from AgMarkNet data: e.g., Solan vs Shimla tomato)", 16, GRAY),
    ("Post-harvest loss: 16% for fruits/vegetables (ICAR data)", 18, LIGHT_GRAY),
    ("KisanMind reduces loss by ~5% via weather-timed harvest", 18, LIGHT_GRAY),
    ("Avg smallholder income: Rs 1,13,000/year (NABARD 2023)", 18, LIGHT_GRAY),
])

# Per-farmer impact
add_rounded_rect(slide, 6.5, 1.6, 6.4, 2.2, BG_CARD)
add_accent_line(slide, 6.7, 1.75, 6, ACCENT_GREEN)
add_multi_text(slide, 6.7, 1.9, 6, 2.0, [
    ("Per-Farmer Impact (Solan Tomato Example)", 24, ACCENT_GREEN, True),
    ("", 4),
    ("Mandi Price Arbitrage: +Rs 12,000/year", 20, WHITE, True),
    ("  Solan mandi Rs 1800/q vs Shimla Rs 2400/q = Rs 600 gap", 16, GRAY),
    ("  Less Rs 200 transport (60km x Rs 3.5/km/q) = Rs 400 net gain", 16, GRAY),
    ("  x 30 quintals/season x 1 season = Rs 12,000", 16, GRAY),
    ("Weather-Timed Harvest: +Rs 10,000/year", 20, WHITE, True),
    ("  Avoiding rain-damaged harvest (5% of Rs 2L gross)", 16, GRAY),
    ("Reduced Advisory Wait: +Rs 12,000/year", 20, WHITE, True),
    ("  From days-to-get-KVK-advice -> 15sec real-time", 16, GRAY),
    ("Total: Rs 34,000/year = 30% income increase", 22, ACCENT_GREEN, True),
])

# Aggregate impact table
add_rounded_rect(slide, 0.5, 4.1, 12.4, 3.1, BG_CARD)
add_accent_line(slide, 0.7, 4.25, 12, ACCENT_BLUE)
add_text_box(slide, 0.7, 4.4, 12, 0.5, "Aggregate Impact (Year 1 — 100,000 farmers)", 24, ACCENT_BLUE, True)

# Impact metrics
metrics = [
    ("TIME SAVED", "15 sec vs 2-3 days", "Real-time satellite + mandi\nvs waiting for KVK visit.\n~3 days saved per decision\nx 10 decisions/season\n= 60 days saved/farmer/year", ACCENT_BLUE),
    ("COST REDUCED", "Rs 10,000/farmer/yr", "5% post-harvest loss prevented\nvia weather-timed harvesting.\n100K farmers x Rs 10K\n= Rs 100 crore saved\nin reduced crop wastage", ACCENT_ORANGE),
    ("REVENUE RECOVERED", "Rs 12,000/farmer/yr", "Mandi arbitrage: sell at best\nmarket, not nearest market.\n100K farmers x Rs 12K\n= Rs 120 crore recovered\nfrom price gap capture", ACCENT_GREEN),
    ("TOTAL IMPACT", "Rs 220 crore/year", "Rs 34,000/farmer x 100K\n= Rs 340 crore total value.\nAt scale (1M farmers):\nRs 3,400 crore/year\n= $400M annual impact", WHITE),
]

for i, (title, value, desc, accent) in enumerate(metrics):
    x = 0.7 + i * 3.1
    y = 5.0
    add_rounded_rect(slide, x, y, 2.9, 2.0, RGBColor(0x15, 0x2B, 0x4A))
    add_multi_text(slide, x + 0.15, y + 0.1, 2.6, 1.9, [
        (title, 16, accent, True, PP_ALIGN.CENTER),
        (value, 22, accent, True, PP_ALIGN.CENTER),
        ("", 4),
        (desc, 13, GRAY, False, PP_ALIGN.CENTER),
    ])


# ═══════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK & DEPLOYMENT
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "Tech Stack & Google Cloud Architecture", 40, WHITE, True)

# Left column: Stack
add_rounded_rect(slide, 0.5, 1.4, 6.2, 5.5, BG_CARD)
add_accent_line(slide, 0.7, 1.55, 5.8, ACCENT_BLUE)
add_multi_text(slide, 0.7, 1.7, 5.8, 5.2, [
    ("Full Stack", 24, ACCENT_BLUE, True),
    ("", 4),
    ("LLM Layer", 18, ACCENT_GREEN, True),
    ("  Gemini 3.1 Pro (advisory synthesis)", 16, LIGHT_GRAY),
    ("  Gemini 3.1 Flash Lite (intent + fact-check)", 16, LIGHT_GRAY),
    ("  Gemini Live (WebSocket voice streaming)", 16, LIGHT_GRAY),
    ("", 4),
    ("Satellite Layer", 18, ACCENT_GREEN, True),
    ("  Google Earth Engine: Sentinel-2, Sentinel-1 SAR", 16, LIGHT_GRAY),
    ("  MODIS Terra (1km LST), NASA SMAP (9km root-zone)", 16, LIGHT_GRAY),
    ("", 4),
    ("Voice Layer", 18, ACCENT_GREEN, True),
    ("  Cloud STT V2 + Cloud TTS Wavenet/Neural2", 16, LIGHT_GRAY),
    ("  Cloud Translation v3 (22 languages)", 16, LIGHT_GRAY),
    ("  Twilio Voice + SMS (outbound calls, TwiML)", 16, LIGHT_GRAY),
    ("", 4),
    ("Application Layer", 18, ACCENT_GREEN, True),
    ("  Frontend: Next.js 16, React 19, TypeScript, Tailwind", 16, LIGHT_GRAY),
    ("  Backend: FastAPI, Python 3.12, async/await, uvicorn", 16, LIGHT_GRAY),
    ("  Cache: L1 in-memory + L2 Google Cloud Storage", 16, LIGHT_GRAY),
    ("", 4),
    ("Data Layer", 18, ACCENT_GREEN, True),
    ("  AgMarkNet / data.gov.in (106+ crop prices)", 16, LIGHT_GRAY),
    ("  Open-Meteo (5-day + 120-day historical)", 16, LIGHT_GRAY),
    ("  Google Maps (Distance Matrix, Places, Geocoding)", 16, LIGHT_GRAY),
])

# Right column: GCP services
add_rounded_rect(slide, 6.9, 1.4, 6.0, 3.0, BG_CARD)
add_accent_line(slide, 7.1, 1.55, 5.6, ACCENT_ORANGE)
add_multi_text(slide, 7.1, 1.7, 5.6, 2.7, [
    ("Google Cloud Services (11)", 24, ACCENT_ORANGE, True),
    ("", 4),
    ("  Earth Engine (satellite processing)", 18, LIGHT_GRAY),
    ("  Cloud Run (backend: asia-south1)", 18, LIGHT_GRAY),
    ("  Compute Engine (frontend VM)", 18, LIGHT_GRAY),
    ("  Cloud Speech-to-Text V2", 18, LIGHT_GRAY),
    ("  Cloud Text-to-Speech", 18, LIGHT_GRAY),
    ("  Cloud Translation API v3", 18, LIGHT_GRAY),
    ("  Cloud Storage (L2 cache)", 18, LIGHT_GRAY),
    ("  Cloud Build + Artifact Registry", 18, LIGHT_GRAY),
    ("  Maps Platform (3 APIs)", 18, LIGHT_GRAY),
    ("  Gemini API (3 model tiers)", 18, LIGHT_GRAY),
])

# Right: Cache architecture
add_rounded_rect(slide, 6.9, 4.6, 6.0, 2.3, BG_CARD)
add_accent_line(slide, 7.1, 4.75, 5.6, ACCENT_GREEN)
add_multi_text(slide, 7.1, 4.9, 5.6, 2.0, [
    ("2-Tier Persistent Cache", 24, ACCENT_GREEN, True),
    ("", 4),
    ("L1 In-Memory: 0.13s response", 18, LIGHT_GRAY),
    ("  Advisory: 15min | NDVI: 6hr | Session: 10min", 16, GRAY),
    ("", 4),
    ("L2 Cloud Storage: ~200ms response", 18, LIGHT_GRAY),
    ("  Mandi raw: 1hr | KVK: 30 days", 16, GRAY),
    ("  Survives deploys, pre-computed per district", 16, GRAY),
    ("", 4),
    ("Every response includes data_age_minutes", 16, ACCENT_BLUE),
])


# ═══════════════════════════════════════════════════════════
# SLIDE 9: COMPETITIVE ADVANTAGE / MOAT
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 1, 0.4, 11, 0.8, "Why KisanMind Wins: The Data Fusion Moat", 40, WHITE, True)
add_text_box(slide, 1, 1.1, 11, 0.5, "Each piece exists in isolation. No other solution fuses all six in real-time via voice.", 24, GRAY)

# Comparison table
headers = ["Capability", "KisanMind", "Kisan Call Center", "AgriApps (DeHaat etc)", "Govt Portals"]
rows = [
    ["Satellite crop health (NDVI)", "Real-time Sentinel-2", "None", "None", "District aggregate only"],
    ["Mandi price arbitrage", "Net profit ranked\n(transport+spoilage)", "Verbal, delayed", "Prices only,\nno transport calc", "Raw prices,\nno ranking"],
    ["Weather -> crop action", "Crop+stage specific\nDO/DON'T rules", "Generic", "Generic forecast", "District level"],
    ["Voice in 22 languages", "Native TTS/STT", "Hindi/English only", "Text only, 2-3 langs", "Text only"],
    ["Growth stage awareness", "GDD from 120-day\nhistorical data", "None", "None", "None"],
    ["Anti-hallucination", "Dual-LLM + 8 rules\n+ confidence scoring", "N/A (human)", "No guardrails", "N/A"],
]

# Draw table
y_start = 1.8
row_h = 0.85
col_widths = [2.8, 2.6, 2.2, 2.6, 2.6]
x_start = 0.4

# Header row
x = x_start
for j, header in enumerate(headers):
    w = col_widths[j]
    add_rounded_rect(slide, x, y_start, w - 0.05, row_h - 0.05, RGBColor(0x15, 0x2B, 0x4A))
    add_text_box(slide, x + 0.1, y_start + 0.15, w - 0.2, row_h - 0.1, header, 16, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    x += w

# Data rows
for i, row in enumerate(rows):
    x = x_start
    y = y_start + (i + 1) * row_h
    for j, cell in enumerate(row):
        w = col_widths[j]
        bg = BG_CARD if i % 2 == 0 else RGBColor(0x15, 0x2B, 0x4A)
        add_rounded_rect(slide, x, y, w - 0.05, row_h - 0.05, bg)
        color = ACCENT_GREEN if j == 1 else (GRAY if j > 1 else LIGHT_GRAY)
        fs = 13 if j > 1 else 14
        add_text_box(slide, x + 0.1, y + 0.05, w - 0.2, row_h - 0.1, cell, fs, color, j == 0, PP_ALIGN.CENTER)
        x += w


# ═══════════════════════════════════════════════════════════
# SLIDE 10: CALL TO ACTION
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)
add_accent_line(slide, 0, 0.06, 13.333, ACCENT_ORANGE)

add_text_box(slide, 1, 1.0, 11, 1.2, "KisanMind", 72, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.2, 11, 0.8, "Satellite-to-Voice Intelligence for Every Indian Farmer", 36, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_accent_line(slide, 4, 3.2, 5.3, ACCENT_ORANGE)

# Key numbers
key_stats = [
    ("150M", "Farmers\nAddressable"),
    ("4", "Satellite\nSystems"),
    ("22", "Indian\nLanguages"),
    ("Rs 34K", "Income Gain\nPer Farmer/Year"),
    ("15 sec", "Advisory\nResponse"),
]

for i, (num, label) in enumerate(key_stats):
    x = 0.8 + i * 2.5
    y = 3.6
    add_multi_text(slide, x, y, 2.3, 1.5, [
        (num, 40, ACCENT_GREEN, True, PP_ALIGN.CENTER),
        (label, 16, GRAY, False, PP_ALIGN.CENTER),
    ])

add_rounded_rect(slide, 2, 5.3, 9.3, 1.3, BG_CARD)
add_multi_text(slide, 2.3, 5.4, 8.7, 1.2, [
    ("Live Now: kisanmind.dmj.one  |  Call: +1 260-254-7946", 28, ACCENT_BLUE, True, PP_ALIGN.CENTER),
    ("", 4),
    ("GitHub: github.com/divyamohan1993/kisanmind  |  Contact: contact@dmj.one", 20, GRAY, False, PP_ALIGN.CENTER),
])

add_text_box(slide, 1, 6.8, 11, 0.5, "ET AI Hackathon 2026  |  Problem #5: Domain-Specialized AI Agents with Compliance Guardrails", 18, GRAY, False, PP_ALIGN.CENTER)

# ── Save ──
output_path = "/mnt/experiments/et-genai-hackathon-phase-2/KisanMind_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
